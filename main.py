import sys
import pysqlite3

# Ensure Chroma uses pysqlite3 (bundled SQLite) in environments where system sqlite is old
sys.modules["sqlite3"] = pysqlite3
sys.modules["sqlite3.dbapi2"] = pysqlite3.dbapi2
 
import os
import re
import io
import uuid
import time
import queue
import json
import hashlib
import mimetypes
import threading
import subprocess
import zipfile
import base64
import tempfile
import torch
from transformers import AutoTokenizer, AutoModelForSequenceClassification
from collections import deque
from pathlib import Path as _Path
from typing import List, Optional, Dict, Any, Tuple, Set

import requests
from fastapi import FastAPI, UploadFile, File, Query, Body, HTTPException, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse, JSONResponse
from pydantic import BaseModel

# ---------------- Vector DB ----------------
import chromadb

# ---------------- File watching ----------------
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# ---------------- Document parsing ----------------
from pypdf import PdfReader
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pdf2image import convert_from_path
import pytesseract
from bs4 import BeautifulSoup
from PIL import Image
import openpyxl
import extract_msg

import traceback
import difflib

# ---------------- Embeddings: nomic-embed-text via ONNXRuntime ----------------
import onnxruntime as ort
from tokenizers import Tokenizer
import numpy as np

class NomicOnnxEmbedder:
    """
    ONNXRuntime wrapper for nomic-ai/nomic-embed-text local model.

    Expected files under EMBED_MODEL_DIR:
      - tokenizer.json
      - onnx/model.onnx   (change onnx_filename if you want a different variant)
    """
    def __init__(self, model_dir: str, max_len: int = None, onnx_filename: str = "model.onnx"):
        import os
        self.model_dir = model_dir.rstrip("/")
        tok_path = os.path.join(self.model_dir, "tokenizer.json")
        onnx_path = os.path.join(self.model_dir, "onnx", onnx_filename)

        if not os.path.exists(tok_path):
            raise RuntimeError(f"Tokenizer not found: {tok_path}")
        if not os.path.exists(onnx_path):
            raise RuntimeError(f"ONNX model not found: {onnx_path}")

        self.tokenizer = Tokenizer.from_file(tok_path)
        self.max_len = max_len or int(os.getenv("EMBED_MAX_LEN", "2048"))

        # CPU by default; switch to onnxruntime-gpu + CUDA providers if desired
        providers = ort.get_available_providers()
        self.session = ort.InferenceSession(onnx_path, providers=providers)

        # Model I/O signatures
        self.input_names = [i.name for i in self.session.get_inputs()]
        outs = [o.name for o in self.session.get_outputs()]
        if not outs:
            raise RuntimeError("Could not resolve ONNX output name.")
        self.output_name = outs[0]  # often 'last_hidden_state'

    def _prepare_arrays(self, texts):
        max_len = self.max_len
        input_ids, attention_mask = [], []
        for t in texts:
            enc = self.tokenizer.encode(t or "")
            ids = enc.ids[:max_len]
            att = [1] * len(ids)
            if len(ids) < max_len:
                pad = max_len - len(ids)
                ids += [0] * pad
                att += [0] * pad
            input_ids.append(ids)
            attention_mask.append(att)

        # Optional inputs: many BERT-style models require these
        token_type_ids = [[0] * max_len for _ in texts]          # all zeros
        position_ids   = [list(range(max_len)) for _ in texts]   # sometimes expected

        arrays = {
            "input_ids":      np.asarray(input_ids, dtype=np.int64),
            "attention_mask": np.asarray(attention_mask, dtype=np.int64),
            "token_type_ids": np.asarray(token_type_ids, dtype=np.int64),
            "position_ids":   np.asarray(position_ids, dtype=np.int64),
        }
        return arrays

    def _masked_mean_pool(self, token_embs: np.ndarray, attn: np.ndarray) -> np.ndarray:
        """
        token_embs: [B, T, H], attn: [B, T]
        returns: [B, H]
        """
        attn = attn[:, :token_embs.shape[1]].astype(np.float32)
        attn_exp = np.expand_dims(attn, axis=-1)                 # [B, T, 1]
        summed = (token_embs * attn_exp).sum(axis=1)             # [B, H]
        counts = np.clip(attn_exp.sum(axis=1), 1e-6, None)       # [B, 1]
        return summed / counts

    def encode(self, texts, normalize_embeddings: bool = True):
        if not isinstance(texts, list):
            texts = [texts]

        arrays = self._prepare_arrays(texts)

        # Feed ONLY what the model actually requires
        feed = {name: arrays[name] for name in self.input_names if name in arrays}
        missing = [name for name in self.input_names if name not in feed]
        if missing:
            raise RuntimeError(f"Missing inputs not covered by adapter: {missing}")

        out = self.session.run([self.output_name], feed)[0]
        vecs = np.asarray(out, dtype=np.float32)

        # Pool if the model returned token embeddings
        if vecs.ndim == 3:                     # [B, T, H]
            vecs = self._masked_mean_pool(vecs, arrays["attention_mask"])
        elif vecs.ndim == 4:
            vecs = vecs.squeeze()
            if vecs.ndim == 3:
                vecs = self._masked_mean_pool(vecs, arrays["attention_mask"])
        elif vecs.ndim == 1:                   # [H] -> [1, H]
            vecs = vecs[None, :]

        # L2 normalize
        if normalize_embeddings:
            norms = np.linalg.norm(vecs, axis=1, keepdims=True) + 1e-12
            vecs = vecs / norms

        return vecs.tolist()



# ---------------- Configuration ----------------
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434")
OLLAMA_KEEP_ALIVE = os.getenv("OLLAMA_KEEP_ALIVE", "720m")
EMBED_MODEL_DIR = os.getenv("EMBED_MODEL_DIR", "/opt/adam/models/nomic-ai/nomic-embed-text")
CHAT_MODEL = os.getenv("CHAT_MODEL", "llama3:8b")  # default to fast 8B
CHROMA_DIR = os.getenv("CHROMA_DIR", "/srv/rag/chroma")
WATCH_DIR = os.getenv("WATCH_DIR", "/srv/rag/watched")
UPLOAD_DIR = os.getenv("UPLOAD_DIR", "/srv/rag/uploads")
COLLECTION = os.getenv("COLLECTION", "company_docs")
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "1400"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "300"))
Q_MAX = int(os.getenv("INGEST_QUEUE_MAX", "8"))
SUMMARY_MODEL = os.getenv("SUMMARY_MODEL", "mistral-7b-instruct")

# ---------------- Model aliasing ----------------
ALIAS_MAP = {
    # Keep any local aliases you still use; avoid mapping removed 70B family
    "Adam Lite": "adam-lite:latest",
    "adam-lite": "adam-lite:latest",
    "llama3:8b": "llama3:8b",
    "mistral-7b-instruct": "mistral-7b-instruct:latest",
}
DEFAULT_MODEL = CHAT_MODEL


def resolve_model(name: Optional[str]) -> str:
    """Return an Ollama model tag given a friendly name or raw tag."""
    if not name:
        return DEFAULT_MODEL
    return ALIAS_MAP.get(name, name)


# Initialize embedder early so startup fails fast if model missing
EMBEDDER = NomicOnnxEmbedder(EMBED_MODEL_DIR)

# ---------------- FastAPI app ----------------
app = FastAPI(title="Local RAG Service")

@app.exception_handler(Exception)
async def all_exception_handler(request, exc):
    tb = "".join(traceback.format_exception(type(exc), exc, exc.__traceback__))
    # Log the full traceback to stdout/journalctl
    print("\n--- Unhandled exception ---\n", tb, flush=True)
    # Return structured JSON so clients don't see plain text "Internal Server Error"
    return JSONResponse(
        status_code=500,
        content={"error": "internal_server_error", "detail": str(exc)},
    )

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

client = chromadb.PersistentClient(path=CHROMA_DIR)
collection = client.get_or_create_collection(COLLECTION)
retriever = collection

_debug_lock = threading.Lock()
_debug_buffer = deque(maxlen=10)
_DEBUG_API_KEY = os.getenv("DEBUG_API_KEY", "")


def _now_ms():
    return int(time.time() * 1000)


def _truncate(s: str, n: int = 280) -> str:
    if s is None:
        return ""
    s = str(s)
    return s if len(s) <= n else s[:n] + "\u2026"


def _get_text_for_hit(hit: dict) -> str:
    """
    Return usable passage text for a retrieved hit.
    Tries multiple common keys (on the hit and in meta/metadata) and falls back.
    """
    for k in ("page_content", "content", "text", "body", "snippet"):
        v = hit.get(k)
        if isinstance(v, str) and v.strip():
            return v.strip()

    meta = hit.get("meta") or hit.get("metadata") or {}
    for k in ("page_content", "content", "text", "body", "snippet"):
        v = meta.get(k)
        if isinstance(v, str) and v.strip():
            return v.strip()

    title = (hit.get("title") or meta.get("title") or "").strip()
    url = (hit.get("sp_web_url") or meta.get("sp_web_url") or hit.get("path") or meta.get("path") or "").strip()
    if title or url:
        return f"{title}\n{url}".strip()

    return ""

# ---- BGE reranker (local/offline) with safe fallback ----
_BGE_PATH = os.getenv("RERANKER_MODEL_PATH", "/opt/rag-models/bge-reranker-v2-m3")
_BGE_DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
_BGE_OK = False
try:
    _BGE_TOK = AutoTokenizer.from_pretrained(_BGE_PATH, local_files_only=True)
    _BGE_MODEL = AutoModelForSequenceClassification.from_pretrained(_BGE_PATH, local_files_only=True)
    _BGE_MODEL.to(_BGE_DEVICE)
    _BGE_MODEL.eval()
    _BGE_OK = True
except Exception:
    _BGE_OK = False


def _bge_scores(query: str, texts: List[str], batch_size: int = 16, max_length: int = 512) -> List[float]:
    scores: List[float] = []
    for i in range(0, len(texts), batch_size):
        bt = texts[i:i+batch_size]
        enc = _BGE_TOK(text=[query]*len(bt), text_pair=bt, truncation=True,
                       max_length=max_length, padding=True, return_tensors="pt")
        enc = {k: v.to(_BGE_DEVICE) for k, v in enc.items()}
        with torch.no_grad():
            logits = _BGE_MODEL(**enc).logits.squeeze(-1)  # [B]
        scores.extend(torch.sigmoid(logits).detach().cpu().tolist())
    return scores


_TOKEN_RE = re.compile(r"\b[\w:.-]{3,}\b", re.UNICODE)


def _tokens_generic(s: str) -> set:
    return set(t.lower() for t in _TOKEN_RE.findall(s or ""))


def _extract_passage(text: str, query: str, window_chars: int = 800) -> str:
    """
    Generic passage picker: prefer segments whose tokens overlap with the query.
    No domain-specific words; purely lexical overlap + small window.
    """
    t = (text or "").strip()
    if not t:
        return ""
    qtok = _tokens_generic(query)
    if not qtok:
        return t[:window_chars]

    # split to sentences (very rough)
    sents = re.split(r'(?<=[.!?])\s+', t)
    # score each sentence by token overlap
    scored = [(i, len(_tokens_generic(s) & qtok), s) for i, s in enumerate(sents)]
    scored.sort(key=lambda x: x[1], reverse=True)
    if not scored or scored[0][1] == 0:
        return t[:window_chars]

    # take a window around the best sentence
    i_best = scored[0][0]
    left = max(0, i_best - 2)
    right = min(len(sents), i_best + 3)
    excerpt = " ".join(sents[left:right]).strip()
    if len(excerpt) < window_chars // 2 and len(t) > window_chars:
        # pad to window size if too short
        start = max(0, t.lower().find(excerpt.lower()) - (window_chars // 4))
        return t[start:start + window_chars]
    return excerpt[:window_chars]

# ---------------- Helpers ----------------
SUPPORTED = (
    ".pdf", ".docx", ".txt",
    ".md", ".csv", ".log", ".json", ".py", ".cs", ".java",
    ".html", ".htm",
    ".eml", ".msg",
    ".xlsx", ".pptx",
    ".png", ".jpg", ".jpeg", ".tiff", ".tif",
    ".doc", ".rtf",
)


def sha1(path: _Path) -> str:
    h = hashlib.sha1()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def read_text(path: _Path) -> str:
    ext = path.suffix.lower()

    if ext == ".pdf":
        # pypdf
        try:
            reader = PdfReader(str(path))
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
            if text and len(text.strip()) > 50:
                return text
        except Exception:
            pass
        # pdftotext
        try:
            out = subprocess.run(["pdftotext", "-layout", str(path), "-"],
                                 capture_output=True, text=True, timeout=60)
            if out.stdout and len(out.stdout.strip()) > 50:
                return out.stdout
        except Exception:
            pass
        # OCR
        try:
            images = convert_from_path(str(path))
            ocr_text = []
            for img in images:
                ocr_text.append(pytesseract.image_to_string(img))
            joined = "\n".join(ocr_text)
            if joined.strip():
                return joined
        except Exception:
            pass
        return ""

    elif ext == ".docx":
        try:
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return ""

    elif ext in (".txt", ".md", ".csv", ".log", ".json", ".py", ".cs", ".java"):
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""

    elif ext in (".html", ".htm"):
        try:
            html = path.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(html, "lxml")
            return soup.get_text(separator="\n")
        except Exception:
            return ""

    elif ext == ".eml":
        try:
            import email
            msg = email.message_from_bytes(path.read_bytes())
            parts = []
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    parts.append(part.get_payload(decode=True).decode(errors="ignore"))
            return "\n".join(parts)
        except Exception:
            return ""

    elif ext == ".msg":
        try:
            m = extract_msg.Message(str(path))
            text = [m.subject or "", m.body or ""]
            return "\n".join([t for t in text if t])
        except Exception:
            return ""

    elif ext == ".xlsx":
        try:
            wb = openpyxl.load_workbook(str(path), data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append(" ".join([str(c) for c in row if c is not None]))
            return "\n".join(parts)
        except Exception:
            return ""

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    if hasattr(shp, "text") and shp.text:
                        parts.append(shp.text)
            return "\n".join(parts)
        except Exception:
            return ""

    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif"):
        try:
            img = Image.open(str(path))
            return pytesseract.image_to_string(img)
        except Exception:
            return ""

    elif ext == ".doc":
        # try antiword
        try:
            out = subprocess.run(["antiword", str(path)], capture_output=True, text=True, timeout=60)
            if out.stdout.strip():
                return out.stdout
        except Exception:
            pass
        # fallback: convert to pdf via libreoffice headless
        try:
            subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                            "--outdir", str(path.parent), str(path)], check=True, timeout=120)
            pdf_path = path.with_suffix(".pdf")
            if pdf_path.exists():
                return read_text(pdf_path)
        except Exception:
            pass
        return ""

    elif ext == ".rtf":
        try:
            out = subprocess.run(["unrtf", "--text", str(path)], capture_output=True, text=True, timeout=60)
            return out.stdout
        except Exception:
            return ""

    else:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""


def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    text = text.strip()
    if not text:
        return []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_text(text)


def embed(texts: List[str]) -> List[List[float]]:
    return EMBEDDER.encode(texts, normalize_embeddings=True)


HEADER_FOOTER_RE = re.compile(r'^(Page\b|Revision\b)', re.IGNORECASE)
SIGNATURE_RE = re.compile(r'\n(?:Regards|Sincerely|Thank you|Thanks|Best)[\s\S]*$', re.IGNORECASE)
LEGAL_RE = re.compile(r'\n(?:Confidentiality Notice|DISCLAIMER:)[\s\S]*$', re.IGNORECASE)

def clean_document_text(text: str) -> str:
    """Basic cleanup: remove headers/footers, signatures, and normalize whitespace."""
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    lines = []
    for ln in text.split('\n'):
        st = ln.strip()
        if HEADER_FOOTER_RE.match(st):
            continue
        lines.append(st)
    cleaned = '\n'.join(lines)
    cleaned = SIGNATURE_RE.sub('', cleaned)
    cleaned = LEGAL_RE.sub('', cleaned)
    cleaned = re.sub(r'\n{2,}', '\n', cleaned)
    cleaned = re.sub(r'[ \t]{2,}', ' ', cleaned)
    return cleaned.strip()


def call_llm(prompt: str) -> str:
    """Send a prompt to the configured LLM and return the raw text response."""
    try:
        r = requests.post(
            f"{OLLAMA_URL}/api/chat",
            json={
                "model": resolve_model(SUMMARY_MODEL),
                "messages": [{"role": "user", "content": prompt}],
                "stream": False,
                "options": {"temperature": 0},
                "keep_alive": OLLAMA_KEEP_ALIVE,
            },
            timeout=120,
        )
        if r.status_code == 200:
            return (r.json().get("message") or {}).get("content", "")
        else:
            print(f"Warning: call_llm HTTP {r.status_code}: {r.text}")
    except Exception as e:
        print(f"Warning: call_llm request failed: {e}")
    return ""


def summarize_document(text_content):
    prompt = f"""
    Summarize the following document content and extract metadata. 
    Respond only with JSON using this exact format:

    {{
      "summary": "<brief summary of the content>",
      "category": "<one or two word category>",
      "keywords": ["<keyword1>", "<keyword2>", ...]
    }}

    Document:
    {text_content}

    Only respond with valid JSON. Do not include any explanatory text or commentary.
    """

    raw_content = call_llm(prompt)
    print("LLM raw content:", raw_content)

    # Try full JSON parse
    try:
        data = json.loads(raw_content)
    except json.JSONDecodeError:
        match = re.search(r'\{.*?\}', raw_content, re.DOTALL)
        if match:
            try:
                data = json.loads(match.group(0))
            except:
                print("Still failed to parse trimmed JSON.")
                data = {}
        else:
            print("Regex match for JSON failed.")
            data = {}

    summary = data.get("summary", "").strip()
    category = data.get("category", "").strip()
    keywords = data.get("keywords", []) if isinstance(data.get("keywords", []), list) else []

    print(f"Parsed summary: {summary}")
    print(f"Parsed category: {category}")
    print(f"Parsed keywords: {keywords}")

    return summary, category, keywords


def upsert_document(path: _Path, source: str) -> int:
    text = read_text(path)
    if not text:
        return 0
    text = clean_document_text(text)
    if len(text) < 500:
        return 0
    summary, category, keywords = summarize_document(text)
    try:
        collection.delete(where={"path": str(path)})
    except Exception:
        pass

    chunks = chunk_text(text)
    if not chunks:
        return 0
    embs = embed(chunks)

    doc_sha = sha1(path)
    ids = [f"{doc_sha}:{i}" for i in range(len(chunks))]
    base_meta = {
        "source": source,
        "path": str(path),
        "summary": summary,
        "category": category,
        "keywords": keywords,
    }
    metas = []
    for i in range(len(chunks)):
        m = _sanitize_metadata(dict(base_meta))
        m["chunk"] = i
        metas.append(m)
    collection.upsert(ids=ids, documents=chunks, metadatas=metas, embeddings=embs)
    return len(chunks)


def _sanitize_metadata(meta: Dict[str, Any]) -> Dict[str, Any]:
    import json
    safe: Dict[str, Any] = {}
    for k, v in (meta or {}).items():
        if isinstance(v, (str, int, float, bool)) or v is None:
            safe[k] = v
        elif isinstance(v, (list, dict, tuple, set)):
            try:
                safe[k] = json.dumps(v, ensure_ascii=False)
            except Exception:
                safe[k] = str(v)
        else:
            safe[k] = str(v)
    return safe



# -------- New: upsert pre-parsed TEXT (SharePoint path) --------
def upsert_text(doc_id: str, text: str, base_meta: Dict[str, Any]) -> int:
    """Upsert pre-parsed TEXT (string) as chunks+embeddings, storing ONLY embeddings+metadata."""
    text = clean_document_text(text or "")
    if len(text) < 500:
        raise ValueError("document under 500 characters after cleanup")

    try:
        collection.delete(where={"doc_id": doc_id})
    except Exception:
        pass
    chunks = chunk_text(text)
    if not chunks:
        return 0
    embs = embed(chunks)
    ids = [f"{doc_id}:{i}" for i in range(len(chunks))]

    # Ensure LLM metadata fields are strings / JSON-serializable
    meta = dict(base_meta or {})
    if not meta.get("summary") and not meta.get("category") and not meta.get("keywords"):
        s, c, k = summarize_document(text)
        meta.update({"summary": s, "category": c, "keywords": k})

    summary = meta.get("summary", "")
    category = meta.get("category", "")
    keywords = meta.get("keywords", [])

    if not isinstance(summary, str):
        try:
            summary = json.dumps(summary, ensure_ascii=False)
        except Exception:
            summary = str(summary)
    if not isinstance(category, str):
        category = str(category)
    if not isinstance(keywords, str):
        try:
            keywords = json.dumps(keywords, ensure_ascii=False)
        except Exception:
            keywords = str(keywords)

    meta.update({"summary": summary, "category": category, "keywords": keywords})
    print(f"upsert_text metadata summary={summary!r}, category={category!r}, keywords={keywords}")

    metas = []
    for i in range(len(chunks)):
        m = dict(meta)
        m["doc_id"] = doc_id
        m["chunk"] = i
        m = _sanitize_metadata(m)
        metas.append(m)
    collection.upsert(ids=ids, documents=chunks, metadatas=metas, embeddings=embs)
    return len(chunks)


def search(query: str, k: int = 4) -> List[Dict[str, Any]]:
    qemb = embed([query])[0]
    res = collection.query(query_embeddings=[qemb], n_results=k,
                           include=["documents", "metadatas", "distances"])
    docs = (res.get("documents") or [[]])[0]
    metas = (res.get("metadatas") or [[]])[0]
    dists = (res.get("distances") or [[]])[0]
    return [{"text": d, "meta": m, "score": float(1.0 / (1e-5 + dist))}
            for d, m, dist in zip(docs, metas, dists)]


def hybrid_rerank(query: str, retriever, reranker_model_name: str,
                  initial_k: int = 15, final_k: int = 5,
                  where: Optional[Dict[str, Any]] = None,
                  debug: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    """Retrieve top N chunks via Chroma then rerank them with a local BGE model."""
    start_retrieve = _now_ms()
    qemb = embed([query])[0]
    query_kwargs = {
        "query_embeddings": [qemb],
        "n_results": initial_k,
        "include": ["documents", "metadatas", "distances"],
    }
    if where:
        query_kwargs["where"] = where
    res = retriever.query(**query_kwargs)
    retrieve_ms = _now_ms() - start_retrieve
    docs = (res.get("documents") or [[]])[0]
    metas = (res.get("metadatas") or [[]])[0]
    dists = (res.get("distances") or [[]])[0]
    chunks = [{"text": d, "meta": m, "score": float(1.0 / (1e-5 + dist))}
              for d, m, dist in zip(docs, metas, dists)]
    if not chunks:
        if debug is not None:
            debug.update({
                "pre_llm_candidates": [],
                "doc_code_boosts": [],
                "reranker": {"raw_text": "", "parsed": [], "threshold": 0.0, "fallback_used": True},
                "post_rerank_hits": [],
                "timing_ms": {"retrieve": retrieve_ms, "rerank": 0},
            })
        return []

    doc_boosts: List[Dict[str, Any]] = []
    _apply_metadata_boost(query, chunks, doc_boosts)

    chunks.sort(key=lambda x: x["score"], reverse=True)
    pre_llm = [dict(ch) for ch in chunks]
    if debug is not None:
        debug["pre_llm_candidates"] = []
        for i, ch in enumerate(pre_llm, 1):
            meta = ch.get("meta") or {}
            debug["pre_llm_candidates"].append({
                "index": i,
                "score_pre": ch.get("score"),
                "doc_code": meta.get("doc_code"),
                "title": meta.get("title"),
                "category": meta.get("category"),
                "revision_date": meta.get("revision_date"),
                "sp_web_url": meta.get("sp_web_url"),
                "snippet": _truncate(ch.get("text") or ch.get("snippet"), 280),
            })
        debug["doc_code_boosts"] = doc_boosts

    # --- BGE rerank attempt over a widened pool ---
    pool = chunks[:max(len(chunks), 40)]  # widen pool (generic)
    texts = [(ch.get("text") or "")[:4000] for ch in pool]

    start_rerank = _now_ms()
    rerank_fallback_used = False
    if _BGE_OK and texts:
        try:
            bge = _bge_scores(query, texts, batch_size=16, max_length=512)
            for ch, s in zip(pool, bge):
                ch["score"] = float(s) * ch.get("_boost", 1.0)
            pool.sort(key=lambda x: x["score"], reverse=True)
            chunks = pool
        except Exception:
            rerank_fallback_used = True
    else:
        rerank_fallback_used = True

    if rerank_fallback_used:
        pass  # fallback to existing order/logic

    result = chunks[:final_k] if chunks else []
    rerank_ms = _now_ms() - start_rerank

    if debug is not None:
        debug["reranker"] = debug.get("reranker", {})
        debug["reranker"].update({
            "raw_text": debug["reranker"].get("raw_text", ""),
            "parsed": [{"index": i + 1, "score": h.get("score")} for i, h in enumerate(result)],
            "threshold": debug["reranker"].get("threshold", 0.0),
            "fallback_used": bool(rerank_fallback_used),
        })
        debug["post_rerank_hits"] = []
        for i, h in enumerate(result, 1):
            meta = h.get("meta") or {}
            debug["post_rerank_hits"].append({
                "index": i,
                "score": h.get("score"),
                "doc_code": meta.get("doc_code"),
                "title": meta.get("title"),
                "snippet": _truncate(h.get("text") or h.get("snippet"), 280),
            })
        debug["timing_ms"] = {"retrieve": retrieve_ms, "rerank": rerank_ms}

    return result




def rerank_sources(question: str, chunks: List[Dict[str, Any]],
                   debug: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    """Rerank retrieved chunks using local BGE scoring."""
    if not chunks:
        if debug is not None:
            debug["reranker"] = {"raw_text": "", "parsed": [], "threshold": 0.0, "fallback_used": True}
        return []

    doc_boosts: List[Dict[str, Any]] = []
    _apply_metadata_boost(question, chunks, doc_boosts)
    texts = []
    for ch in chunks:
        t = _get_text_for_hit(ch)
        texts.append(t[:4000])
    fallback_used = False
    if _BGE_OK:
        try:
            scores = _bge_scores(question, texts)
            for ch, s in zip(chunks, scores):
                ch["score"] = float(s) * ch.get("_boost", 1.0)
        except Exception as e:
            print("Rerank failed:", e)
            fallback_used = True
            for ch in chunks:
                ch["score"] = ch.get("score", 0.0)
    else:
        fallback_used = True

    chunks.sort(key=lambda x: x["score"], reverse=True)
    top = chunks[:3]

    if debug is not None:
        debug.update({
            "doc_code_boosts": doc_boosts,
            "reranker": {
                "raw_text": "",
                "parsed": [{"index": i + 1, "score": h.get("score")} for i, h in enumerate(top)],
                "threshold": 0.0,
                "fallback_used": fallback_used,
            },
        })
    return top


def _dehedge(text: str) -> str:
    patterns = [
        r'(?i)^\s*(according to|based on|from)\s+(the\s+)?(provided|given)\s+(context|information|documents)\s*[:,\-]*\s*',
        r'(?i)^\s*(according to|based on)\s+(the\s+)?document(s)?\s*[:,\-]*\s*',
    ]
    for p in patterns:
        text = re.sub(p, "", text, count=1)
    return text.strip()

def ask_with_context(question: str, hits: List[dict], chat_history: Optional[List[dict]] = None,
                     model: Optional[str] = None, force_citations: bool = False,
                     extra_system_prompt: str = "") -> str:
    ql = (question or "").lower()

    meta_triggers = [
        "your name", "what is your name", "what's your name",
        "who are you", "who is adam", "what does adam stand for",
        "what are you", "introduce yourself"
    ]
    if any(t in ql for t in meta_triggers):
        return "I am Adam - the Amentum Document and Assistance Model (ADAM)."

    context = "\n\n".join([f"[{h.get('index', i+1)}] {h['text']}" for i, h in enumerate(hits)])

    sys_prompt = (
        "You are Adam — the Amentum Document and Assistance Model (ADAM). "
        "Answer directly and succinctly. Do not start with phrases like "
        "'According to the provided context'. Use ONLY the provided context for factual claims and insert "
        "inline bracket citations like [1], [2] right after the sentence they support. "
        "Do not append a 'Sources:' section. If the answer is not in the context, say you do not know. "
        "CITATION RULES: Use only numeric bracket citations that correspond to the provided context blocks, e.g., [1], [2]. "
        "Do not use section numbers like [4.1], ranges like [1-3], or textual citations. Put the citation immediately after each claim it supports. "
        "COMPLETENESS RULES: If the user asks about definitions, boundaries, windows, or procedures, include all relevant elements present in the context (e.g., start and end times, total hours, tool/system names). "
        "NO FABRICATION: If a claim cannot be supported with a bracket citation from the context, say you do not know. "
        "OUTPUT FORMAT (HTML ONLY): Respond with a well-formed HTML fragment (not a full <html> page). Use <p> for paragraphs (each sentence starts with a capital letter). <ul> / <ol> with <li> for lists of steps or bullets. <table><thead>…</thead><tbody>…</tbody></table> for side-by-side facts. <strong>, <em>, <code>, <sup> as needed. CITATIONS: Put bracket citations inline at the end of the clause they support as <sup>[n]</sup>. Only use numbers that map to the provided context blocks. STYLE & SAFETY: Do not include <script>, inline CSS, external images, or arbitrary attributes. No markdown; HTML only. COMPLETENESS: For definition/boundary/steps questions, include all relevant elements present in the context (e.g., start and end times, total hours, and system/tool names). "
        "Direct answer in one sentence. One short follow-up sentence with any missing critical detail (e.g., 'begins Friday 12:00 noon and ends next Friday 11:59 a.m.'). Include the bracket citations inline."
    )
    if force_citations:
        sys_prompt += " You MUST include at least one citation if you answer. If unsure, say you do not know."
    if extra_system_prompt:
        sys_prompt += " " + extra_system_prompt

    messages = [{"role": "system", "content": sys_prompt}]
    if chat_history:
        messages.extend(chat_history)
    prompt = f"Context:\n{context}\n\nQuestion: {question}\nAnswer:"
    messages.append({"role": "user", "content": prompt})

    model_tag = resolve_model(model)

    try:
        r = requests.post(
            f"{OLLAMA_URL}/api/chat",
            json={
                "model": model_tag,
                "messages": messages,
                "stream": False,
                "options": {"temperature": 0.2, "num_predict": 1024, "top_p": 0.95},
                "keep_alive": OLLAMA_KEEP_ALIVE  # keep the 8B resident
            },
            timeout=120
        )
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Ollama connection error: {e}")

    if r.status_code != 200:
        body = r.text
        if len(body) > 800:
            body = body[:800] + "...(truncated)"
        raise HTTPException(status_code=502, detail=f"Ollama status {r.status_code}: {body}")

    try:
        j = r.json()
    except Exception:
        body = r.text
        if len(body) > 800:
            body = body[:800] + "...(truncated)"
        raise HTTPException(status_code=502, detail=f"Ollama returned non-JSON: {body}")

    msg = j.get("message") or {}
    content = msg.get("content") or j.get("content") or ""

    if not content.strip() and hits:
        top = (hits[0].get("text") or "").strip()
        if top and len(top) <= 200 and any(p in ql for p in ["what does", "say", "content", "quote", "exact text"]):
            return f'It says: "{{ " ".join(top.split()) }}" [1].'

    return _dehedge(content)


def filter_cited_sources(answer: str, chunks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Return chunks that are explicitly cited in the answer."""
    if not answer or not chunks:
        return []
    used = {int(n) for n in re.findall(r"\[(\d+)\]", answer) if n.isdigit()}
    if not used:
        return []
    ordered = []
    for ch in chunks:
        if ch.get("index") in used:
            ordered.append(ch)
    return ordered


def _extract_numeric_citations(answer: str, max_index: int) -> Tuple[Set[int], bool]:
    brackets = re.findall(r"\[([^\]]+)\]", answer or "")
    used: Set[int] = set()
    valid = True
    for b in brackets:
        if not b.isdigit():
            valid = False
            continue
        idx = int(b)
        if idx < 1 or idx > max_index:
            valid = False
        else:
            used.add(idx)
    return used, valid


def _ensure_html(text: str) -> Tuple[str, bool, List[str]]:
    tags = [t.lower() for t in re.findall(r"<\s*([a-zA-Z0-9]+)", text or "")]
    is_html = bool(tags)
    if not any(t in ("p", "ul", "ol", "table") for t in tags):
        t = (text or "").strip()
        if t:
            t = t[0].upper() + t[1:]
        text = f"<p>{t}</p>"
        tags = ["p"]
        is_html = True
    tag_summary = sorted({t for t in tags})
    return text, is_html, tag_summary


def _ensure_doc_code(meta: Dict[str, Any]) -> None:
    if meta.get("doc_code"):
        return
    url = meta.get("sp_web_url") or meta.get("path") or ""
    m = re.search(r"([A-Za-z]{2,}-[A-Za-z]{2,}-[A-Za-z]{2,}-\d{3,5})", url)
    if m:
        meta["doc_code"] = m.group(1)


def _apply_metadata_boost(query: str, chunks: List[Dict[str, Any]], doc_boosts: Optional[List[Dict[str, Any]]] = None) -> None:
    def _tokens(s: str) -> List[str]:
        return re.findall(r"\b\w+\b", (s or "").lower())

    q_tokens = set(_tokens(query))
    code_pat = re.compile(r"[A-Za-z0-9]+(?:-[A-Za-z0-9]+)+")
    query_codes = {c.lower() for c in code_pat.findall(query)}
    if doc_boosts is None:
        doc_boosts = []
    for idx, ch in enumerate(chunks, 1):
        meta = ch.get("meta") or {}
        _ensure_doc_code(meta)
        fields = []
        cat = meta.get("category")
        if cat:
            fields.append(cat)
        title = meta.get("title")
        if title:
            fields.append(title)
        kws = meta.get("keywords")
        if kws:
            if isinstance(kws, list):
                fields.extend(kws)
            else:
                fields.append(kws)
        doc_code = meta.get("doc_code")
        if doc_code:
            fields.append(doc_code)
        version = meta.get("version")
        if version:
            fields.append(str(version))
        rev_date = meta.get("revision_date")
        if rev_date:
            fields.append(str(rev_date))

        meta_tokens = set()
        for f in fields:
            meta_tokens.update(_tokens(str(f)))

        matches = 0
        for qt in q_tokens:
            if qt in meta_tokens:
                matches += 1
            else:
                for mt in meta_tokens:
                    if difflib.SequenceMatcher(None, qt, mt).ratio() >= 0.8:
                        matches += 1
                        break

        boost = 1.0
        if matches:
            boost *= 1 + 0.1 * matches
        if doc_code:
            dc = str(doc_code).lower()
            for qc in query_codes:
                if dc == qc:
                    boost *= 3.0
                    doc_boosts.append({"index": idx, "reason": "exact", "multiplier": 3.0})
                    break
                norm_dc = dc.replace("-", "")
                norm_qc = qc.replace("-", "")
                if norm_dc == norm_qc or difflib.SequenceMatcher(None, norm_dc, norm_qc).ratio() >= 0.9:
                    boost *= 2.0
                    doc_boosts.append({"index": idx, "reason": "fuzzy", "multiplier": 2.0})
                    break
        if (meta.get("category") or "").lower() == "policy":
            boost *= 1.15
        ch.setdefault("_boost", boost)
        ch["score"] *= boost

def rewrite_prompt(prompt: str) -> str:
    """Use Mistral to safely rewrite vague prompts, without changing intent."""
    try:
        code_pat = re.compile(r"[A-Za-z0-9]+(?:-[A-Za-z0-9]+)+")
        codes = code_pat.findall(prompt)

        messages = [
            {
                "role": "system",
                "content": (
                    "You are a helpful assistant that rewrites vague document questions into clearer, more precise ones "
                    "without changing their meaning. Do NOT guess or introduce new topics. "
                    "Only rewrite if the original query is unclear or incomplete."
                )
            },
            {
                "role": "user",
                "content": f"Original query: {prompt}\n\nRewritten query:",
            }
        ]

        r = requests.post(
            f"{OLLAMA_URL}/api/chat",
            json={
                "model": "mistral-7b-instruct",
                "messages": messages,
                "stream": False,
                "options": {"temperature": 0.3},
                "keep_alive": OLLAMA_KEEP_ALIVE
            },
            timeout=30
        )

        if r.status_code == 200:
            j = r.json()
            rewritten = (j.get("message") or {}).get("content", "").strip()
            if rewritten:
                for c in codes:
                    if c not in rewritten:
                        rewritten = (rewritten + " " + c).strip()
                orig_lower = prompt.lower()
                rew_lower = rewritten.lower()
                keywords = ["workweek", "work week", "pto", "decisions tool", "timekeeping"]
                for kw in keywords:
                    if kw in orig_lower and kw not in rew_lower:
                        rewritten = (rewritten + " " + kw).strip()
                        rew_lower = rewritten.lower()
                if not codes or all(c in rewritten for c in codes):
                    return rewritten

        return prompt  # fallback if empty or invalid rewrite
    except Exception as e:
        print("Prompt rewrite failed:", e)
        return prompt

# ---- API models & endpoints ----
class QueryBody(BaseModel):
    query: str
    k: int = 4
    history: Optional[List[dict]] = None
    model: Optional[str] = None  # "Adam Large", "Adam Lite", or raw Ollama tag
    # New optional filters (non-breaking)
    org: Optional[str] = None
    category: Optional[str] = None
    doc_code: Optional[str] = None
    owner: Optional[str] = None
    rewrite: bool = True


class QueryResponse(BaseModel):
    """Response schema for the /query endpoint."""
    answer: str
    sources: List[Dict[str, Any]]
    original_query: str
    rewritten_query: str
    prompt_rewritten: bool


# Quick smoke tests:
# curl -s -X POST http://localhost:8000/query -H 'Content-Type: application/json' \
#   -d '{"query": "When does the standard workweek begin and end?"}'
# curl -s -X POST http://localhost:8000/query -H 'Content-Type: application/json' \
#   -d '{"query": "Per CLG-EN-PO-0301, where do I submit PTO?"}'
# curl -s -X POST http://localhost:8000/query -H 'Content-Type: application/json' \
#   -d '{"query": "What is the weather on Mars?"}'

@app.post("/query", response_model=QueryResponse)
def query_api(body: QueryBody) -> QueryResponse:
    start_total = _now_ms()
    code_pat = re.compile(r"[A-Za-z0-9]+(?:-[A-Za-z0-9]+)+")
    debug = {
        "request_id": str(uuid.uuid4()),
        "timestamp_ms": _now_ms(),
        "endpoint": "/query",
        "model": {"generator": resolve_model(body.model), "reranker": "bge-reranker-v2-m3 (transformers, local)" if _BGE_OK else "mistral-7b-instruct"},
        "query": {
            "original": body.query,
            "rewritten": "",
            "prompt_rewritten": False,
            "doc_codes_original": code_pat.findall(body.query),
            "doc_codes_rewritten": [],
        },
        "retrieval": {"where": None},
        "answering": {
            "first_pass": {"answer": "", "citations": [], "timing_ms": 0},
            "second_pass_invoked": False,
            "second_pass": {"answer": "", "citations": [], "timing_ms": 0},
            "final": {
                "answer_used": "",
                "filtered_sources": [],
                "fallback_reason": "none",
                "answer_is_html": False,
                "html_tag_summary": [],
            },
        },
        "total_time_ms": 0,
    }

    if body.rewrite:
        rewritten_query_raw = rewrite_prompt(body.query).strip()
        m = re.search(r"Rewritten Query:\s*(.*)", rewritten_query_raw, re.DOTALL)
        if m:
            rewritten_query = m.group(1).strip()
        else:
            rewritten_query = next(
                (ln.strip() for ln in rewritten_query_raw.splitlines() if ln.strip()),
                body.query,
            )
        prompt_rewritten = rewritten_query != body.query
    else:
        rewritten_query = body.query
        prompt_rewritten = False
    debug["query"]["rewritten"] = rewritten_query
    debug["query"]["prompt_rewritten"] = prompt_rewritten
    debug["query"]["doc_codes_rewritten"] = code_pat.findall(rewritten_query)

    where: Optional[Dict[str, Any]] = None
    if body.org or body.category or body.doc_code or body.owner:
        where = {}
        if body.org:
            where["org"] = body.org
        if body.category:
            where["category"] = body.category
        if body.doc_code:
            where["doc_code"] = body.doc_code
        if body.owner:
            where["owner"] = body.owner
    debug["retrieval"]["where"] = where

    rdebug: Dict[str, Any] = {}
    top_hits = hybrid_rerank(rewritten_query, retriever, "bge-reranker-v2-m3", where=where, debug=rdebug)
    debug["retrieval"].update(rdebug)
    if not top_hits:
        final_answer = "I'm sorry, I couldn't find relevant information."
        final_answer, is_html_final, tags_final = _ensure_html(final_answer)
        debug["answering"]["final"].update({
            "answer_used": _truncate(final_answer, 1200),
            "filtered_sources": [],
            "fallback_reason": "no_hits_after_rerank",
            "answer_is_html": is_html_final,
            "html_tag_summary": tags_final,
        })
        debug["total_time_ms"] = _now_ms() - start_total
        with _debug_lock:
            _debug_buffer.append(debug)
        return QueryResponse(
            answer=final_answer,
            sources=[],
            original_query=body.query,
            rewritten_query=rewritten_query,
            prompt_rewritten=prompt_rewritten,
        )
    hits: List[Dict[str, Any]] = []
    for h in top_hits:
        t = _get_text_for_hit(h)
        if t:
            h["text"] = t
            hits.append(h)
    if not hits:
        for h in top_hits:
            t = _get_text_for_hit(h)
            if t:
                h["text"] = t
                hits.append(h)
            if len(hits) >= 3:
                break

    if not hits:
        final_answer = (
            "<p>I couldn't load readable text from the retrieved sources. "
            "Please re-index the policy with a small text preview.</p>"
        )
        debug["answering"]["final"].update({
            "answer_used": _truncate(final_answer, 1200),
            "filtered_sources": [],
            "fallback_reason": "no_readable_text",
            "answer_is_html": True,
            "html_tag_summary": ["p"],
        })
        debug["total_time_ms"] = _now_ms() - start_total
        with _debug_lock:
            _debug_buffer.append(debug)
        return QueryResponse(
            answer=final_answer,
            sources=[],
            original_query=body.query,
            rewritten_query=rewritten_query,
            prompt_rewritten=prompt_rewritten,
        )

    for i, h in enumerate(hits):
        raw = h.get("text") or ""
        h["text"] = _extract_passage(raw, rewritten_query)
        h["index"] = i + 1

    debug.setdefault("retrieval", {}).setdefault("context_block_summaries", [
        {
            "index": i + 1,
            "chars": len(h.get("text") or ""),
            "preview": (h.get("text") or "")[:120] + ("\u2026" if len(h.get("text") or "") > 120 else ""),
        }
        for i, h in enumerate(hits)
    ])

    start_ans = _now_ms()
    answer_first = ask_with_context(rewritten_query, hits, chat_history=body.history, model=body.model)
    ans_time = _now_ms() - start_ans
    answer_first, is_html_first, tags_first = _ensure_html(answer_first)
    used_first, valid_first = _extract_numeric_citations(answer_first, len(hits))
    citations_first = sorted(list(used_first))
    debug["answering"]["first_pass"] = {
        "answer": _truncate(answer_first, 1200),
        "citations": citations_first,
        "timing_ms": ans_time,
    }
    filtered = filter_cited_sources(answer_first, hits) if valid_first and used_first else []
    answer_used = answer_first
    is_html_final = is_html_first
    tags_final = tags_first
    if not valid_first or not filtered:
        debug["answering"]["second_pass_invoked"] = True
        extra = (
            "Your previous answer used invalid citations (e.g., [4.1] or out-of-range). "
            "Rewrite the answer using only numeric citations [1..N] corresponding to the provided context blocks. "
            "Rewrite the answer as HTML per the output rules above, and fix citations to numeric [1..N]. "
            "Maintain completeness and sentence capitalization."
        )
        snippet_all = " ".join([h.get("text") or "" for h in hits])
        if "12:00" in snippet_all and "11:59" in snippet_all:
            extra += " Include both the start and end time and the total hours if mentioned."
        start_second = _now_ms()
        second = ask_with_context(
            rewritten_query,
            hits,
            chat_history=body.history,
            model=body.model,
            force_citations=True,
            extra_system_prompt=extra,
        )
        second_ms = _now_ms() - start_second
        second, is_html_second, tags_second = _ensure_html(second)
        used_second, valid_second = _extract_numeric_citations(second, len(hits))
        citations_second = sorted(list(used_second))
        debug["answering"]["second_pass"] = {
            "answer": _truncate(second, 1200),
            "citations": citations_second,
            "timing_ms": second_ms,
        }
        if valid_second and citations_second:
            filtered = filter_cited_sources(second, hits)
            if filtered:
                answer_used = second
                is_html_final = is_html_second
                tags_final = tags_second
        if not filtered:
            final_answer = "I'm sorry, I can't answer confidently from the provided sources."
            final_answer, is_html_final, tags_final = _ensure_html(final_answer)
            debug["answering"]["final"].update({
                "answer_used": _truncate(final_answer, 1200),
                "filtered_sources": [],
                "fallback_reason": "no_citations_after_second_pass",
                "answer_is_html": is_html_final,
                "html_tag_summary": tags_final,
            })
            debug["total_time_ms"] = _now_ms() - start_total
            with _debug_lock:
                _debug_buffer.append(debug)
            return QueryResponse(
                answer=final_answer,
                sources=[],
                original_query=body.query,
                rewritten_query=rewritten_query,
                prompt_rewritten=prompt_rewritten,
            )
    else:
        debug["answering"]["second_pass"] = {"answer": "", "citations": [], "timing_ms": 0}

    rich = []
    for h in filtered:
        meta = h.get("meta", {}) or {}
        rich.append({
            "index": h.get("index"),
            "title": meta.get("title"),
            "org": meta.get("org"),
            "category": meta.get("category"),
            "version": meta.get("version"),
            "revision_date": meta.get("revision_date"),
            "doc_code": meta.get("doc_code"),
            "sp_web_url": meta.get("sp_web_url"),
            "path": meta.get("path"),   # present for legacy local docs
            "score": h.get("score"),
            "snippet": _get_text_for_hit(h)[:280],
        })

    debug["answering"]["final"].update({
        "answer_used": _truncate(answer_used, 1200),
        "filtered_sources": [
            {
                "index": s["index"],
                "doc_code": s["doc_code"],
                "title": s["title"],
                "sp_web_url": s["sp_web_url"],
                "score": s["score"],
                "snippet": _truncate(s["snippet"], 280),
            }
            for s in rich
        ],
        "answer_is_html": is_html_final,
        "html_tag_summary": tags_final,
    })
    debug["total_time_ms"] = _now_ms() - start_total
    with _debug_lock:
        _debug_buffer.append(debug)

    return QueryResponse(
        answer=answer_used,
        sources=rich,
        original_query=body.query,
        rewritten_query=rewritten_query,
        prompt_rewritten=prompt_rewritten,
    )


@app.post("/upload")
def upload_api(file: UploadFile = File(...)):
    dest = _Path(UPLOAD_DIR) / file.filename
    dest.parent.mkdir(parents=True, exist_ok=True)
    with dest.open("wb") as f:
        f.write(file.file.read())
    n = upsert_document(dest, source="upload")
    return {"ok": True, "chunks": n, "path": str(dest)}


@app.post("/reindex")
def reindex_api():
    root = _Path(WATCH_DIR)
    count, total_chunks = 0, 0
    if not root.exists():
        return {"files": 0, "chunks": 0}
    for path in root.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED:
            n = upsert_document(path, source="watched")
            total_chunks += n
            count += 1
    return {"files": count, "chunks": total_chunks}


@app.get("/status")
def status():
    try:
        cnt = collection.count()
    except Exception:
        cnt = None
    return {"ok": True, "collection": COLLECTION, "count": cnt, "supported": list(SUPPORTED)}


@app.delete("/delete")
def delete_by_path(path: str = Query(..., description="absolute file path stored in metadata")):
    try:
        collection.delete(where={"path": path})
        return {"ok": True, "path": path}
    except Exception as e:
        return {"ok": False, "error": str(e)}


@app.post("/search")
def search_only(payload: dict = Body(...)):
    q = payload.get("query", "")
    k = int(payload.get("k", 4))
    return {"results": search(q, k=k)}


@app.get("/list")
def list_docs(limit: int = 10000):
    data = collection.get(include=["metadatas"], limit=limit)
    paths = sorted({(m or {}).get("path", "unknown") for m in data.get("metadatas", [])})
    return {"count": len(paths), "paths": paths}


@app.post("/reset")
def reset_api():
    try:
        client.delete_collection(COLLECTION)
    except Exception:
        pass
    global collection
    collection = client.get_or_create_collection(COLLECTION)
    return {"ok": True}

@app.get("/ollama_health")
def ollama_health():
    try:
        r = requests.get(f"{OLLAMA_URL}/api/tags", timeout=10)
        ct = r.headers.get("content-type", "")
        body = r.json() if "application/json" in ct else r.text
        return {"ok": r.status_code == 200, "status": r.status_code, "body": body}
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Ollama unreachable: {e}")

@app.get("/embed_health")
def embed_health():
    try:
        vec = EMBEDDER.encode(["hello world"])[0]
        return {"ok": True, "dim": len(vec), "preview": vec[:8]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Embed error: {e}")


@app.get("/list_docs")
def list_documents():
    try:
        collection = collection = client.get_or_create_collection(COLLECTION)

        # Fetch all document entries with metadata
        results = collection.get(include=["metadatas", "documents"], limit=10000)

        docs = []
        for i in range(len(results["ids"])):
            meta = results["metadatas"][i] or {}
            doc_info = {
                "doc_id": results["ids"][i],
                "metadata": meta,
            }
         
         
            docs.append(doc_info)

        return JSONResponse(content={"documents": docs, "count": len(docs)})

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


# ---- folder watcher ----
class Handler(FileSystemEventHandler):
    def on_created(self, event):
        if event.is_directory:
            return
        p = _Path(event.src_path)
        if p.suffix.lower() in SUPPORTED:
            try:
                upsert_document(p, source="watched")
            except Exception as e:
                print("ingest error:", e)

    def on_modified(self, event):
        self.on_created(event)


if _Path(WATCH_DIR).exists():
    obs = Observer()
    obs.schedule(Handler(), WATCH_DIR, recursive=True)
    obs.start()


# ---------------- Ingest queue & jobs ----------------
INGEST_Q: "queue.Queue[tuple[str,str,str]]" = queue.Queue(maxsize=Q_MAX)
JOBS: dict[str, dict] = {}


def _ingest_worker():
    while True:
        job_id, path, source = INGEST_Q.get()
        try:
            t0 = time.time()
            chunks = upsert_document(_Path(path), source=source)
            JOBS[job_id] = {"status": "done", "path": path, "chunks": chunks, "ms": int((time.time() - t0) * 1000)}
        except Exception as e:
            JOBS[job_id] = {"status": "error", "path": path, "error": str(e)}
        finally:
            INGEST_Q.task_done()


threading.Thread(target=_ingest_worker, daemon=True).start()


@app.post("/upload_async")
def upload_async_api(file: UploadFile = File(...)):
    dest = _Path(UPLOAD_DIR) / file.filename
    dest.parent.mkdir(parents=True, exist_ok=True)
    with dest.open("wb") as f:
        f.write(file.file.read())
    jid = str(uuid.uuid4())
    JOBS[jid] = {"status": "queued", "path": str(dest)}
    try:
        INGEST_Q.put_nowait((jid, str(dest), "upload"))
    except queue.Full:
        JOBS[jid] = {"status": "error", "path": str(dest), "error": "ingest queue is full"}
    return {"ok": True, "job_id": jid, "path": str(dest)}


@app.get("/jobs")
def list_jobs():
    keys = list(JOBS.keys())[-200:]
    return {k: JOBS[k] for k in keys}


@app.get("/jobs/{job_id}")
def job_status(job_id: str):
    return JOBS.get(job_id, {"status": "unknown"})


# ---------------- Downloads ----------------
def _is_under_allowed(p: _Path) -> bool:
    roots = [_Path(UPLOAD_DIR).resolve(), _Path(WATCH_DIR).resolve()]
    try:
        rp = p.resolve()
        for r in roots:
            try:
                if rp.is_relative_to(r.resolve()):
                    return True
            except AttributeError:
                # Python 3.9 doesn't have Path.is_relative_to
                if str(rp).startswith(str(r.resolve()) + os.sep):
                    return True
        return False
    except Exception:
        return False


@app.get("/download")
def download(path: str, inline: bool = False):
    p = _Path(path)
    if not p.exists() or not p.is_file():
        raise HTTPException(404, "file not found")
    if not _is_under_allowed(p):
        raise HTTPException(403, "forbidden path")

    mt, _ = mimetypes.guess_type(str(p))
    resp = FileResponse(
        path=str(p),
        media_type=mt or "application/octet-stream",
        filename=p.name,
    )
    if inline:
        resp.headers["Content-Disposition"] = f'inline; filename="{p.name}"'
    return resp


class ZipRequest(BaseModel):
    paths: List[str]


@app.post("/download_zip")
def download_zip(req: ZipRequest):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for raw in req.paths or []:
            p = _Path(raw)
            if not (p.exists() and p.is_file() and _is_under_allowed(p)):
                continue
            z.write(str(p), arcname=p.name)
    buf.seek(0)
    headers = {"Content-Disposition": 'attachment; filename="documents.zip"'}
    return StreamingResponse(buf, media_type="application/zip", headers=headers)


# ---------------- Filtered search ----------------
def search_filtered(
    query: str,
    k: int = 4,
    path: Optional[str] = None,
    paths: Optional[List[str]] = None,
    prefix: Optional[str] = None,
    source: Optional[str] = None,
    org: Optional[str] = None,
    category: Optional[str] = None,
    doc_code: Optional[str] = None,
    owner: Optional[str] = None,
) -> List[Dict[str, Any]]:
    qemb = embed([query])[0]
    include = ["documents", "metadatas", "distances"]

    def to_hits(res):
        docs = (res.get("documents") or [[]])[0]
        metas = (res.get("metadatas") or [[]])[0]
        dists = (res.get("distances") or [[]])[0]
        return [{"text": d, "meta": m, "score": float(1.0 / (1e-5 + dist))}
                for d, m, dist in zip(docs, metas, dists)]

    results: List[Dict[str, Any]] = []

    if path:
        res = collection.query(query_embeddings=[qemb], n_results=max(k, 8),
                               include=include, where={"path": path})
        results.extend(to_hits(res))
    elif paths:
        for p in paths:
            res = collection.query(query_embeddings=[qemb], n_results=max(2, k // max(1, len(paths))) + 6,
                                   include=include, where={"path": p})
            results.extend(to_hits(res))
    elif prefix:
        try:
            res = collection.query(query_embeddings=[qemb], n_results=max(k, 12),
                                   include=include, where={"path": {"$contains": prefix}})
            results.extend(to_hits(res))
        except Exception:
            res = collection.query(query_embeddings=[qemb], n_results=max(k * 5, 20), include=include)
            results.extend([h for h in to_hits(res) if str(h["meta"].get("path", "")).startswith(prefix)])
    elif source or org or category or doc_code or owner:
        where: Dict[str, Any] = {}
        if source:   where["source"] = source
        if org:      where["org"] = org
        if category: where["category"] = category
        if doc_code: where["doc_code"] = doc_code
        if owner:    where["owner"] = owner
        res = collection.query(query_embeddings=[qemb], n_results=max(k, 12),
                               include=include, where=where)
        results.extend(to_hits(res))
    else:
        res = collection.query(query_embeddings=[qemb], n_results=k, include=include)
        results.extend(to_hits(res))

    seen: Dict[Tuple[Any, Any], Dict[str, Any]] = {}
    for h in results:
        key = (h["meta"].get("doc_id") or h["meta"].get("path"), h["meta"].get("chunk"))
        if key not in seen or h["score"] > seen[key]["score"]:
            seen[key] = h
    hits = sorted(seen.values(), key=lambda x: x["score"], reverse=True)[:k]
    return hits


@app.post("/query_path")
def query_path(body: Dict[str, Any] = Body(...)):
    start_total = _now_ms()
    query = body.get("query", "")
    k = int(body.get("k", 4))
    history = body.get("history")
    path = body.get("path")
    paths = body.get("paths")
    prefix = body.get("prefix")
    source = body.get("source")
    model = body.get("model")

    org = body.get("org")
    category = body.get("category")
    doc_code = body.get("doc_code")
    owner = body.get("owner")

    code_pat = re.compile(r"[A-Za-z0-9]+(?:-[A-Za-z0-9]+)+")
    where_filters = {}
    if path:
        where_filters["path"] = path
    if paths:
        where_filters["paths"] = paths
    if prefix:
        where_filters["prefix"] = prefix
    if source:
        where_filters["source"] = source
    if org:
        where_filters["org"] = org
    if category:
        where_filters["category"] = category
    if doc_code:
        where_filters["doc_code"] = doc_code
    if owner:
        where_filters["owner"] = owner

    debug = {
        "request_id": str(uuid.uuid4()),
        "timestamp_ms": _now_ms(),
        "endpoint": "/query_path",
        "model": {"generator": resolve_model(model), "reranker": "bge-reranker-v2-m3 (transformers, local)" if _BGE_OK else "mistral-7b-instruct"},
        "query": {
            "original": query,
            "rewritten": query,
            "prompt_rewritten": False,
            "doc_codes_original": code_pat.findall(query),
            "doc_codes_rewritten": code_pat.findall(query),
        },
        "retrieval": {"where": where_filters or None,
                       "pre_llm_candidates": [],
                       "doc_code_boosts": [],
                       "reranker": {"raw_text": "", "parsed": [], "threshold": 0.0, "fallback_used": False},
                       "post_rerank_hits": [],
                       "timing_ms": {"retrieve": 0, "rerank": 0}},
        "answering": {
            "first_pass": {"answer": "", "citations": [], "timing_ms": 0},
            "second_pass_invoked": False,
            "second_pass": {"answer": "", "citations": [], "timing_ms": 0},
            "final": {
                "answer_used": "",
                "filtered_sources": [],
                "fallback_reason": "none",
                "answer_is_html": False,
                "html_tag_summary": [],
            },
        },
        "total_time_ms": 0,
    }

    start_retrieve = _now_ms()
    hits0 = search_filtered(query, k=10, path=path, paths=paths, prefix=prefix, source=source,
                            org=org, category=category, doc_code=doc_code, owner=owner)
    retrieve_ms = _now_ms() - start_retrieve
    for i, ch in enumerate(hits0, 1):
        meta = ch.get("meta") or {}
        _ensure_doc_code(meta)
        debug["retrieval"]["pre_llm_candidates"].append({
            "index": i,
            "score_pre": ch.get("score"),
            "doc_code": meta.get("doc_code"),
            "title": meta.get("title"),
            "category": meta.get("category"),
            "revision_date": meta.get("revision_date"),
            "sp_web_url": meta.get("sp_web_url"),
            "snippet": _truncate(ch.get("text") or ch.get("snippet"), 280),
        })

    rdebug: Dict[str, Any] = {}
    start_rerank = _now_ms()
    top_hits = rerank_sources(query, hits0, debug=rdebug)
    rerank_ms = _now_ms() - start_rerank
    debug["retrieval"].update(rdebug)
    debug["retrieval"]["timing_ms"] = {"retrieve": retrieve_ms, "rerank": rerank_ms}
    debug["retrieval"]["post_rerank_hits"] = [
        {
            "index": i + 1,
            "score": h.get("score"),
            "doc_code": (h.get("meta") or {}).get("doc_code"),
            "title": (h.get("meta") or {}).get("title"),
            "snippet": _truncate(_get_text_for_hit(h), 280),
        }
        for i, h in enumerate(top_hits)
    ]

    if not top_hits:
        final_answer = "I'm sorry, I couldn't find relevant information."
        final_answer, is_html_final, tags_final = _ensure_html(final_answer)
        debug["answering"]["final"].update({
            "answer_used": _truncate(final_answer, 1200),
            "filtered_sources": [],
            "fallback_reason": "no_hits_after_rerank",
            "answer_is_html": is_html_final,
            "html_tag_summary": tags_final,
        })
        debug["total_time_ms"] = _now_ms() - start_total
        with _debug_lock:
            _debug_buffer.append(debug)
        return {"answer": final_answer, "sources": []}

    hits: List[Dict[str, Any]] = []
    for h in top_hits:
        t = _get_text_for_hit(h)
        if t:
            h["text"] = t
            hits.append(h)
    if not hits:
        for h in top_hits:
            t = _get_text_for_hit(h)
            if t:
                h["text"] = t
                hits.append(h)
            if len(hits) >= 3:
                break

    if not hits:
        final_answer = (
            "<p>I couldn't load readable text from the retrieved sources. Please re-index the policy with a small text preview.</p>"
        )
        debug["answering"]["final"].update({
            "answer_used": _truncate(final_answer, 1200),
            "filtered_sources": [],
            "fallback_reason": "no_readable_text",
            "answer_is_html": True,
            "html_tag_summary": ["p"],
        })
        debug["total_time_ms"] = _now_ms() - start_total
        with _debug_lock:
            _debug_buffer.append(debug)
        return {"answer": final_answer, "sources": []}

    for i, h in enumerate(hits):
        raw = h.get("text") or ""
        h["text"] = _extract_passage(raw, query)
        h["index"] = i + 1

    debug.setdefault("retrieval", {}).setdefault("context_block_summaries", [
        {
            "index": i + 1,
            "chars": len(h.get("text") or ""),
            "preview": (h.get("text") or "")[:120] + ("\u2026" if len(h.get("text") or "") > 120 else ""),
        }
        for i, h in enumerate(hits)
    ])

    start_ans = _now_ms()
    answer_first = ask_with_context(query, hits, chat_history=history, model=model)
    ans_ms = _now_ms() - start_ans
    answer_first, is_html_first, tags_first = _ensure_html(answer_first)
    used_first, valid_first = _extract_numeric_citations(answer_first, len(hits))
    citations_first = sorted(list(used_first))
    debug["answering"]["first_pass"] = {
        "answer": _truncate(answer_first, 1200),
        "citations": citations_first,
        "timing_ms": ans_ms,
    }
    filtered = filter_cited_sources(answer_first, hits) if valid_first and used_first else []
    answer_used = answer_first
    is_html_final = is_html_first
    tags_final = tags_first
    if not valid_first or not filtered:
        debug["answering"]["second_pass_invoked"] = True
        extra = (
            "Your previous answer used invalid citations (e.g., [4.1] or out-of-range). "
            "Rewrite the answer using only numeric citations [1..N] corresponding to the provided context blocks. "
            "Rewrite the answer as HTML per the output rules above, and fix citations to numeric [1..N]. "
            "Maintain completeness and sentence capitalization."
        )
        snippet_all = " ".join([h.get("text") or "" for h in hits])
        if "12:00" in snippet_all and "11:59" in snippet_all:
            extra += " Include both the start and end time and the total hours if mentioned."
        start_second = _now_ms()
        second = ask_with_context(
            query,
            hits,
            chat_history=history,
            model=model,
            force_citations=True,
            extra_system_prompt=extra,
        )
        second_ms = _now_ms() - start_second
        second, is_html_second, tags_second = _ensure_html(second)
        used_second, valid_second = _extract_numeric_citations(second, len(hits))
        citations_second = sorted(list(used_second))
        debug["answering"]["second_pass"] = {
            "answer": _truncate(second, 1200),
            "citations": citations_second,
            "timing_ms": second_ms,
        }
        if valid_second and citations_second:
            filtered = filter_cited_sources(second, hits)
            if filtered:
                answer_used = second
                is_html_final = is_html_second
                tags_final = tags_second
        if not filtered:
            final_answer = "I'm sorry, I can't answer confidently from the provided sources."
            final_answer, is_html_final, tags_final = _ensure_html(final_answer)
            debug["answering"]["final"].update({
                "answer_used": _truncate(final_answer, 1200),
                "filtered_sources": [],
                "fallback_reason": "no_citations_after_second_pass",
                "answer_is_html": is_html_final,
                "html_tag_summary": tags_final,
            })
            debug["total_time_ms"] = _now_ms() - start_total
            with _debug_lock:
                _debug_buffer.append(debug)
            return {"answer": final_answer, "sources": []}
    else:
        debug["answering"]["second_pass"] = {"answer": "", "citations": [], "timing_ms": 0}
    for h in filtered:
        h["snippet"] = _get_text_for_hit(h)[:280]

    debug["answering"]["final"].update({
        "answer_used": _truncate(answer_used, 1200),
        "filtered_sources": [
            {
                "index": h.get("index"),
                "doc_code": (h.get("meta") or {}).get("doc_code"),
                "title": (h.get("meta") or {}).get("title"),
                "sp_web_url": (h.get("meta") or {}).get("sp_web_url"),
                "score": h.get("score"),
                "snippet": _truncate(h.get("snippet"), 280),
            }
            for h in filtered
        ],
        "answer_is_html": is_html_final,
        "html_tag_summary": tags_final,
    })
    debug["total_time_ms"] = _now_ms() - start_total
    with _debug_lock:
        _debug_buffer.append(debug)

    return {"answer": answer_used, "sources": filtered}


# Latest record (without snippets)
# curl -s http://localhost:8000/debug_last | jq .
#
# Latest with snippets (truncate to 400 chars)
# curl -s "http://localhost:8000/debug_last?include_snippets=true&max_chars=400" | jq .
#
# Third most recent
# curl -s "http://localhost:8000/debug_last?index=-3" | jq .
#
# With auth (if DEBUG_API_KEY=secret123 is set)
# curl -s -H "X-Debug-Key: secret123" "http://localhost:8000/debug_last?include_snippets=true" | jq .


@app.get("/debug_last")
def debug_last(index: int = -1,
               include_snippets: bool = False,
               max_chars: int = 280,
               x_debug_key: Optional[str] = Header(default=None)):
    if _DEBUG_API_KEY and x_debug_key != _DEBUG_API_KEY:
        raise HTTPException(status_code=401, detail="Unauthorized")
    count = len(_debug_buffer)
    if count == 0:
        return {"count": 0, "index": None, "record": None}
    if index < 0:
        idx = count + index
    else:
        idx = index
    if idx < 0 or idx >= count:
        raise HTTPException(status_code=400, detail="index out of range")
    record = _debug_buffer[idx]
    rec = json.loads(json.dumps(record))

    def sanitize(obj):
        if isinstance(obj, dict):
            out = {}
            for k, v in obj.items():
                if isinstance(v, (dict, list)):
                    out[k] = sanitize(v)
                elif isinstance(v, str):
                    if include_snippets:
                        out[k] = _truncate(v, max_chars)
                    else:
                        out[k] = "" if k in ("snippet", "answer", "raw_text", "answer_used") else v
                else:
                    out[k] = v
            return out
        elif isinstance(obj, list):
            return [sanitize(x) for x in obj]
        elif isinstance(obj, str):
            return _truncate(obj, max_chars) if include_snippets else obj
        else:
            return obj

    rec = sanitize(rec)
    return {"count": count, "index": idx, "record": rec}


@app.post("/debug_last/clear")
def debug_last_clear(x_debug_key: Optional[str] = Header(default=None)):
    if _DEBUG_API_KEY and x_debug_key != _DEBUG_API_KEY:
        raise HTTPException(status_code=401, detail="Unauthorized")
    with _debug_lock:
        cnt = len(_debug_buffer)
        _debug_buffer.clear()
    return {"cleared": cnt}


# ---------------- New: Pydantic model for SharePoint ingest ----------------
class IngestDocument(BaseModel):
    # Identity & links (from SharePoint)
    sp_site_id: Optional[str] = None
    sp_list_id: Optional[str] = None
    sp_item_id: Optional[str] = None
    sp_drive_id: Optional[str] = None
    sp_file_id: Optional[str] = None
    sp_web_url: str

    # Versioning
    etag: Optional[str] = None
    version_label: Optional[str] = None

    # Core metadata
    title: Optional[str] = None
    doc_code: Optional[str] = None
    org_code: Optional[str] = None
    org: Optional[str] = None
    category: Optional[str] = None
    owner: Optional[str] = None
    version: Optional[str] = None
    revision_date: Optional[str] = None
    latest_review_date: Optional[str] = None
    document_review_date: Optional[str] = None
    review_approval_date: Optional[str] = None
    keywords: Optional[List[str]] = None
    enterprise_keywords: Optional[List[str]] = None
    association_ids: Optional[List[str]] = None
    domain: Optional[str] = "HR"
    allowed_groups: Optional[List[str]] = None

    # Content options (you’ll use content_bytes)
    file_name: Optional[str] = None
    content_bytes: Optional[str] = None   # base64 of the file bytes
    text_content: Optional[str] = None    # if you pre-extract text in your SP worker

    # Optional overrides
    chunk_size: Optional[int] = None
    chunk_overlap: Optional[int] = None
    persist: Optional[bool] = False       # ignored; never persist files locally


# ---------------- New: SharePoint-first ingest (no local persistence) ----------------
@app.post("/ingest_document")
def ingest_document_api(body: IngestDocument):
    """
    Accept SharePoint metadata + content (base64 bytes or pre-extracted text),
    parse to text in-memory (if needed), chunk, embed, upsert. No file persistence.
    """
    # Use per-request chunk overrides if provided
    global CHUNK_SIZE, CHUNK_OVERLAP
    orig_size, orig_overlap = CHUNK_SIZE, CHUNK_OVERLAP
    try:
        if body.chunk_size: CHUNK_SIZE = int(body.chunk_size)
        if body.chunk_overlap: CHUNK_OVERLAP = int(body.chunk_overlap)

        # Build stable versioned doc_id
        version_key = body.version_label or body.etag or "v1"
        spid = str(body.sp_item_id or body.sp_file_id or body.sp_web_url)
        doc_id = f"{spid}:{version_key}"

        # Base metadata attached to every chunk
        base_meta = {
            "source": "sharepoint",
            "sp_site_id": body.sp_site_id,
            "sp_list_id": body.sp_list_id,
            "sp_item_id": body.sp_item_id,
            "sp_drive_id": body.sp_drive_id,
            "sp_file_id": body.sp_file_id,
            "sp_web_url": body.sp_web_url,
            "etag": body.etag,
            "version_label": body.version_label,
            "summary" : "",
            "title": body.title,
            "doc_code": body.doc_code,
            "org_code": body.org_code,
            "org": body.org,
            "category": body.category,
            "owner": body.owner,
            "version": body.version,
            "revision_date": body.revision_date,
            "latest_review_date": body.latest_review_date,
            "document_review_date": body.document_review_date,
            "review_approval_date": body.review_approval_date,
            "keywords": body.keywords or [],
            "enterprise_keywords": body.enterprise_keywords or [],
            "association_ids": body.association_ids or [],
            "domain": body.domain or "HR",
            "allowed_groups": body.allowed_groups or ["AllEmployees"],
        }

        # 1) If text provided, use directly (fastest)
        if body.text_content and str(body.text_content).strip():
            text = clean_document_text(body.text_content)
            if len(text) < 500:
                raise HTTPException(status_code=400, detail="document under 500 characters after cleanup")
            summary, category, keywords  = summarize_document(text)
            meta = dict(base_meta)
            meta.update({"summary": summary, "category": category, "keywords": keywords})
            n = upsert_text(doc_id, text, meta)
            return {"ok": True, "doc_id": doc_id, "chunks": n, "used": "text_content", "summary":summary}

        # 2) Else parse from base64 content bytes ephemerally
        if body.content_bytes:
            raw = base64.b64decode(body.content_bytes)
            suffix = _Path(body.file_name).suffix if body.file_name else ""
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
                tmp = _Path(tf.name)
                tf.write(raw)
            try:
                text = read_text(tmp)
            finally:
                try:
                    tmp.unlink(missing_ok=True)
                except TypeError:
                    if tmp.exists():
                        tmp.unlink()
            text = clean_document_text(text)
            if len(text) < 500:
                raise HTTPException(status_code=400, detail="document under 500 characters after cleanup")
            summary, category, keywords = summarize_document(text)
            meta = dict(base_meta)
            meta.update({"summary": summary, "category": category, "keywords": keywords})
            n = upsert_text(doc_id, text, meta)
            return {"ok": True, "doc_id": doc_id, "chunks": n, "used": "content_bytes"}

        raise HTTPException(status_code=400, detail="Provide either text_content or content_bytes.")
    finally:
        CHUNK_SIZE, CHUNK_OVERLAP = orig_size, orig_overlap


# ---------------- New: Delete by SharePoint identity ----------------
@app.delete("/delete_sp")
def delete_by_sp(sp_item_id: Optional[str] = Query(None),
                 sp_file_id: Optional[str] = Query(None),
                 version: Optional[str] = Query(None),
                 etag: Optional[str] = Query(None)):
    """
    Delete all chunks for a SharePoint document identity. If version/etag is
    supplied, only that version is removed; otherwise all versions are removed.
    """
    where: Dict[str, Any] = {}
    if sp_item_id:
        where["sp_item_id"] = sp_item_id
    if sp_file_id:
        where["sp_file_id"] = sp_file_id
    if not where:
        raise HTTPException(400, "Provide sp_item_id or sp_file_id")
    if version:
        where["version_label"] = version
    if etag:
        where["etag"] = etag
    try:
        collection.delete(where=where)
        return {"ok": True, "where": where}
    except Exception as e:
        return {"ok": False, "error": str(e), "where": where}
