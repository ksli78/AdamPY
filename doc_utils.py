import hashlib
import subprocess
from pathlib import Path as _Path
from typing import Any, Dict, List
import json

from bs4 import BeautifulSoup
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from PIL import Image
from pypdf import PdfReader
from pdf2image import convert_from_path
import pytesseract
import openpyxl
import extract_msg



def sha1(path: _Path) -> str:
    h = hashlib.sha1()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def read_text(path: _Path) -> str:
    ext = path.suffix.lower()

    if ext == ".pdf":
        # pypdf
        try:
            reader = PdfReader(str(path))
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
            if text and len(text.strip()) > 50:
                return text
        except Exception:
            pass
        # pdftotext
        try:
            import subprocess
            out = subprocess.run(["pdftotext", str(path), "-"] , capture_output=True, text=True, timeout=60)
            return out.stdout
        except Exception:
            pass
        # pdf2image + OCR
        try:
            images = convert_from_path(str(path))
            return "\n".join(pytesseract.image_to_string(img) for img in images)
        except Exception:
            return ""

    elif ext == ".docx":
        try:
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return ""

    elif ext in (".txt", ".md", ".csv", ".log", ".json", ".py", ".cs", ".java"):
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""

    elif ext in (".html", ".htm"):
        try:
            html = path.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(html, "lxml")
            return soup.get_text(separator="\n")
        except Exception:
            return ""

    elif ext == ".eml":
        try:
            import email
            msg = email.message_from_bytes(path.read_bytes())
            parts = []
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    parts.append(part.get_payload(decode=True).decode(errors="ignore"))
            return "\n".join(parts)
        except Exception:
            return ""

    elif ext == ".msg":
        try:
            m = extract_msg.Message(str(path))
            text = [m.subject or "", m.body or ""]
            return "\n".join([t for t in text if t])
        except Exception:
            return ""

    elif ext == ".xlsx":
        try:
            wb = openpyxl.load_workbook(str(path), data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append(" ".join([str(c) for c in row if c is not None]))
            return "\n".join(parts)
        except Exception:
            return ""

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    if hasattr(shp, "text") and shp.text:
                        parts.append(shp.text)
            return "\n".join(parts)
        except Exception:
            return ""

    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif"):
        try:
            img = Image.open(str(path))
            return pytesseract.image_to_string(img)
        except Exception:
            return ""

    elif ext == ".doc":
        # try antiword
        try:
            out = subprocess.run(["antiword", str(path)], capture_output=True, text=True, timeout=60)
            if out.stdout.strip():
                return out.stdout
        except Exception:
            pass
        # fallback: convert to pdf via libreoffice headless
        try:
            subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                            "--outdir", str(path.parent), str(path)], check=True, timeout=120)
            pdf_path = path.with_suffix(".pdf")
            if pdf_path.exists():
                return read_text(pdf_path)
        except Exception:
            pass
        return ""

    elif ext == ".rtf":
        try:
            out = subprocess.run(["unrtf", "--text", str(path)], capture_output=True, text=True, timeout=60)
            return out.stdout
        except Exception:
            return ""

    else:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""


def chunk_text(text: str, size: int = None, overlap: int = None) -> List[str]:
    from config import CHUNK_SIZE, CHUNK_OVERLAP
    size = size or CHUNK_SIZE
    overlap = overlap or CHUNK_OVERLAP
    text = text.strip()
    if not text:
        return []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_text(text)


def sanitize_metadata(meta: Dict[str, Any]) -> Dict[str, Any]:
    safe: Dict[str, Any] = {}
    for k, v in (meta or {}).items():
        if isinstance(v, (str, int, float, bool)) or v is None:
            safe[k] = v
        elif isinstance(v, (list, dict, tuple, set)):
            try:
                safe[k] = json.dumps(v, ensure_ascii=False)
            except Exception:
                safe[k] = str(v)
        else:
            safe[k] = str(v)
    return safe
